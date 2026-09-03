"""PPT generation service — SVG-based pipeline using ppt-master's svg_to_pptx engine.

Pipeline:
  1. LLM generates a JSON outline (title, subtitle, slides[{title, bullets, notes}])
     plus a design spec (colors, fonts, layout hints)
  2. Each slide is rendered as an SVG file (1280 × 960, ppt169)
  3. ppt-master's create_pptx_with_native_svg converts SVG → PPTX
"""
from __future__ import annotations

import json
import logging
import os
import re
import sys
import tempfile
import textwrap
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.core.config import settings

logger = logging.getLogger(__name__)

# ── ppt-master path (OPTIONAL enhancement) ──────────────────────────────
# ppt-master is an optional skill that enables pixel-perfect SVG → PPTX import.
# If it is not present (e.g. after a fresh clone without the large skill dir),
# we transparently fall back to the built-in python-pptx renderer, so PPT
# export always works with zero extra setup.
_PPTMASTER_CANDIDATES = [
    Path(__file__).resolve().parents[2] / "skills" / "ppt-master" / "scripts",   # backend/skills/...
    Path(__file__).resolve().parents[3] / "skills" / "ppt-master" / "scripts",   # repo-root skills/...
    Path(__file__).resolve().parent / "ppt_master",                              # vendored app/services/ppt_master
]
_PPTMASTER_SCRIPTS = next((p for p in _PPTMASTER_CANDIDATES if p.exists()), None)
_PPTMASTER_AVAILABLE = _PPTMASTER_SCRIPTS is not None

if _PPTMASTER_AVAILABLE and str(_PPTMASTER_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_PPTMASTER_SCRIPTS))

# ── Canvas: ppt169 = 1280 × 960 px ─────────────────────────────────────
CANVAS_W = 1280
CANVAS_H = 960
CANVAS_FORMAT = "ppt169"

# ─────────────────────────────────────────────────────────────────────────
#  SVG SYSTEM PROMPT
# ─────────────────────────────────────────────────────────────────────────

SVG_SYSTEM_PROMPT = textwrap.dedent("""
You are a professional slide designer. You will receive a JSON outline with a global design spec
and a list of slides. For each slide, produce ONE SVG file (1280×960) that looks great.

### SVG Technical Rules (mandatory — violations break PPTX export)
- viewBox MUST be "0 0 1280 960"
- NO <style>, NO class attributes, ALL styles inline (fill=, stroke=, font-family=, etc.)
- NO <mask>, NO <symbol>/<use>, NO <foreignObject>, NO <script>, NO @font-face
- NO HTML named entities (&nbsp; &mdash; &copy; etc.) — use raw Unicode characters instead
- XML reserved characters MUST be escaped: & → &amp;  < → &lt;  > → &gt;
- Text in Chinese/English mixed: use font-family="Arial, Microsoft YaHei, sans-serif"
- NO <animate*>, NO event attributes, NO <iframe>

### Design Guidelines
- Title slide: large centered title, subtitle below, decorative background elements
- Content slides: clear hierarchy — slide title at top, main content below
- Use the colors/fonts from the design spec
- Vary layouts: bullet list, two-column, icon+text, stat cards, timeline steps, etc.
- Keep text readable: min font-size 18px for body, 28px for titles
- Add visual interest: colored rectangles, circles, lines as design elements
- DO NOT use images (<image> tags) unless given a data URI

### Output Format
Return a JSON array of SVG strings (one per slide), in this exact format:
[
  "<svg viewBox=\\"0 0 1280 960\\" xmlns=\\"http://www.w3.org/2000/svg\\">...</svg>",
  "<svg viewBox=\\"0 0 1280 960\\" xmlns=\\"http://www.w3.org/2000/svg\\">...</svg>"
]
Output ONLY the JSON array. No markdown fences. No explanation.
""").strip()


# ─────────────────────────────────────────────────────────────────────────
#  OUTLINE PROMPT
# ─────────────────────────────────────────────────────────────────────────

OUTLINE_SYSTEM_PROMPT = textwrap.dedent("""
You are a professional presentation strategist. Analyze the source text and produce a
structured presentation outline in JSON format.

Output language: {language}

User style description: {style_description}

Design a color scheme and font plan matching the style description. Return JSON in this structure:
{{
  "title": "Presentation Title",
  "subtitle": "Subtitle or author/date",
  "design": {{
    "bg": "#FFFFFF",
    "primary": "#1A56E8",
    "accent": "#F59E0B",
    "text": "#1F2937",
    "text_light": "#6B7280",
    "font_title": "Arial",
    "font_body": "Arial",
    "dark_mode": false
  }},
  "slides": [
    {{
      "title": "Slide Title",
      "layout": "bullets | two_col | stats | timeline | quote | section_break",
      "bullets": ["point 1", "point 2", "point 3"],
      "notes": "Speaker notes for this slide"
    }}
  ]
}}

Guidelines:
- 8-14 slides total (including title slide as slide 0 automatically)
- Use varied layout types for visual interest
- Each slide: 3-5 bullet points max
- First slide in the "slides" array is the first content slide (title slide is generated separately)
- Include a conclusion/summary slide
- Output ONLY valid JSON, no markdown fences
""").strip()


# ─────────────────────────────────────────────────────────────────────────
#  SVG BUILDERS (fast fallback — no LLM)
# ─────────────────────────────────────────────────────────────────────────

def _esc(text: str) -> str:
    """Escape XML reserved characters."""
    return (str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))


def _hex_luminance(hex_color: str) -> float:
    """Return relative luminance (0–1) of a hex color."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return 0.5
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrasting_text(bg: str, light: str = "#FFFFFF", dark: str = "#1F2937") -> str:
    return light if _hex_luminance(bg) < 0.45 else dark


def _build_title_svg(title: str, subtitle: str, design: dict) -> str:
    bg = design.get("bg", "#1A56E8")
    primary = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    text_color = _contrasting_text(bg, "#FFFFFF", "#1F2937")
    font = design.get("font_title", "Arial")

    # Decorative circle
    circle_fill = accent

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <!-- Decorative background circles -->
  <circle cx="1150" cy="120" r="260" fill="{_esc(primary)}" opacity="0.15"/>
  <circle cx="-80" cy="860" r="300" fill="{_esc(accent)}" opacity="0.10"/>
  <!-- Left accent bar -->
  <rect x="0" y="0" width="8" height="960" fill="{_esc(primary)}"/>
  <!-- Bottom accent bar -->
  <rect x="0" y="920" width="1280" height="40" fill="{_esc(primary)}" opacity="0.8"/>
  <!-- Title -->
  <text x="640" y="400" text-anchor="middle" font-family="{_esc(font)}, Microsoft YaHei, sans-serif"
        font-size="72" font-weight="bold" fill="{_esc(text_color)}">{_esc(title)}</text>
  <!-- Accent line -->
  <rect x="440" y="430" width="400" height="5" fill="{_esc(accent)}" rx="2"/>
  <!-- Subtitle -->
  <text x="640" y="510" text-anchor="middle" font-family="{_esc(font)}, Microsoft YaHei, sans-serif"
        font-size="32" fill="{_esc(text_color)}" opacity="0.75">{_esc(subtitle)}</text>
  <!-- Decorative dot -->
  <circle cx="640" cy="590" r="8" fill="{_esc(accent)}"/>
</svg>'''


def _build_bullets_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("bg", "#FFFFFF")
    primary = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    text = design.get("text", "#1F2937")
    text_light = design.get("text_light", "#6B7280")
    font_title = design.get("font_title", "Arial")
    font_body = design.get("font_body", "Arial")

    title = slide.get("title", "")
    bullets = slide.get("bullets", [])

    # Alternate subtle tint on even/odd slides
    tint = "#F8FAFF" if idx % 2 == 0 else bg

    bullet_svgs = []
    for i, b in enumerate(bullets[:6]):
        y = 290 + i * 90
        bullet_svgs.append(
            f'  <circle cx="96" cy="{y - 8}" r="10" fill="{_esc(primary)}"/>'
        )
        bullet_svgs.append(
            f'  <text x="124" y="{y}" font-family="{_esc(font_body)}, Microsoft YaHei, sans-serif"'
            f' font-size="32" fill="{_esc(text)}">{_esc(b)}</text>'
        )

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(tint)}"/>
  <!-- Top accent bar -->
  <rect x="0" y="0" width="1280" height="8" fill="{_esc(primary)}"/>
  <!-- Header background -->
  <rect x="0" y="8" width="1280" height="140" fill="{_esc(primary)}"/>
  <!-- Slide number badge -->
  <rect x="1160" y="28" width="80" height="48" rx="8" fill="{_esc(accent)}"/>
  <text x="1200" y="61" text-anchor="middle" font-family="{_esc(font_title)}, sans-serif"
        font-size="24" font-weight="bold" fill="#FFFFFF">{idx + 1}</text>
  <!-- Title -->
  <text x="64" y="100" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="44" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>
  <!-- Bullets -->
{chr(10).join(bullet_svgs)}
  <!-- Bottom line -->
  <rect x="0" y="900" width="1280" height="4" fill="{_esc(accent)}" opacity="0.6"/>
</svg>'''


def _build_section_break_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    font_title = design.get("font_title", "Arial")
    title = slide.get("title", "")
    bullets = slide.get("bullets", [])
    subtitle = bullets[0] if bullets else ""

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <circle cx="1100" cy="200" r="350" fill="#FFFFFF" opacity="0.05"/>
  <circle cx="180" cy="800" r="250" fill="{_esc(accent)}" opacity="0.10"/>
  <rect x="0" y="920" width="1280" height="40" fill="{_esc(accent)}"/>
  <text x="640" y="420" text-anchor="middle" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="80" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>
  <rect x="440" y="450" width="400" height="5" fill="{_esc(accent)}" rx="2"/>
  <text x="640" y="530" text-anchor="middle" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="36" fill="#FFFFFF" opacity="0.8">{_esc(subtitle)}</text>
</svg>'''


def _build_two_col_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("bg", "#FFFFFF")
    primary = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    text = design.get("text", "#1F2937")
    font_title = design.get("font_title", "Arial")
    font_body = design.get("font_body", "Arial")

    title = slide.get("title", "")
    bullets = slide.get("bullets", [])
    mid = len(bullets) // 2 or 1
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    def _col_items(items: list, x: int) -> list[str]:
        lines = []
        for i, b in enumerate(items[:4]):
            y = 340 + i * 100
            lines.append(
                f'  <rect x="{x}" y="{y - 28}" width="6" height="38" rx="3" fill="{_esc(accent)}"/>'
            )
            lines.append(
                f'  <text x="{x + 24}" y="{y}" font-family="{_esc(font_body)}, Microsoft YaHei, sans-serif"'
                f' font-size="30" fill="{_esc(text)}">{_esc(b)}</text>'
            )
        return lines

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <rect x="0" y="0" width="1280" height="8" fill="{_esc(primary)}"/>
  <rect x="0" y="8" width="1280" height="140" fill="{_esc(primary)}"/>
  <text x="64" y="100" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="44" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>
  <!-- Divider -->
  <rect x="630" y="200" width="4" height="700" fill="{_esc(primary)}" opacity="0.15"/>
  <!-- Left column header -->
  <text x="80" y="240" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="28" font-weight="bold" fill="{_esc(primary)}">Key Points</text>
  <!-- Right column header -->
  <text x="680" y="240" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="28" font-weight="bold" fill="{_esc(primary)}">Details</text>
{chr(10).join(_col_items(left_bullets, 80))}
{chr(10).join(_col_items(right_bullets, 680))}
  <rect x="0" y="900" width="1280" height="4" fill="{_esc(accent)}" opacity="0.6"/>
</svg>'''


def _build_stats_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("bg", "#FFFFFF")
    primary = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    text = design.get("text", "#1F2937")
    font_title = design.get("font_title", "Arial")
    font_body = design.get("font_body", "Arial")

    title = slide.get("title", "")
    bullets = slide.get("bullets", [])
    cards = bullets[:4]
    n = len(cards)
    card_w = 260
    total_w = n * card_w + (n - 1) * 30
    start_x = (1280 - total_w) // 2

    card_svgs = []
    for i, b in enumerate(cards):
        cx = start_x + i * (card_w + 30)
        card_svgs.extend([
            f'  <rect x="{cx}" y="320" width="{card_w}" height="320" rx="16" fill="{_esc(primary)}"/>',
            f'  <text x="{cx + card_w//2}" y="460" text-anchor="middle"'
            f' font-family="{_esc(font_body)}, Microsoft YaHei, sans-serif"'
            f' font-size="26" fill="#FFFFFF">{_esc(b)}</text>',
        ])

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <rect x="0" y="0" width="1280" height="8" fill="{_esc(primary)}"/>
  <rect x="0" y="8" width="1280" height="140" fill="{_esc(primary)}"/>
  <text x="64" y="100" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="44" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>
{chr(10).join(card_svgs)}
  <rect x="0" y="900" width="1280" height="4" fill="{_esc(accent)}" opacity="0.6"/>
</svg>'''


def _build_quote_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    font_title = design.get("font_title", "Arial")

    title = slide.get("title", "")
    bullets = slide.get("bullets", [])
    quote = bullets[0] if bullets else ""
    source = bullets[1] if len(bullets) > 1 else ""

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <circle cx="200" cy="200" r="300" fill="#FFFFFF" opacity="0.04"/>
  <!-- Large quote mark -->
  <text x="100" y="380" font-family="Georgia, serif" font-size="200" fill="{_esc(accent)}" opacity="0.5">&quot;</text>
  <!-- Quote text -->
  <text x="640" y="380" text-anchor="middle" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="40" font-style="italic" fill="#FFFFFF">{_esc(quote)}</text>
  <!-- Source -->
  <text x="640" y="480" text-anchor="middle" font-family="{_esc(font_title)}, sans-serif"
        font-size="28" fill="{_esc(accent)}">{_esc(source)}</text>
  <!-- Title -->
  <text x="64" y="840" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="32" fill="#FFFFFF" opacity="0.7">{_esc(title)}</text>
  <rect x="0" y="920" width="1280" height="40" fill="{_esc(accent)}" opacity="0.8"/>
</svg>'''


def _build_timeline_svg(slide: dict, design: dict, idx: int) -> str:
    bg = design.get("bg", "#FFFFFF")
    primary = design.get("primary", "#1A56E8")
    accent = design.get("accent", "#F59E0B")
    text = design.get("text", "#1F2937")
    font_title = design.get("font_title", "Arial")
    font_body = design.get("font_body", "Arial")

    title = slide.get("title", "")
    bullets = slide.get("bullets", [])

    steps = bullets[:5]
    n = len(steps)
    if n == 0:
        return _build_bullets_svg(slide, design, idx)

    step_w = 1100 // max(n, 1)
    start_x = 90

    step_svgs = []
    for i, s in enumerate(steps):
        cx = start_x + i * step_w + step_w // 2
        # Connector line (not for last)
        if i < n - 1:
            step_svgs.append(
                f'  <line x1="{cx + 40}" y1="480" x2="{cx + step_w - 40}" y2="480"'
                f' stroke="{_esc(primary)}" stroke-width="3" opacity="0.4"/>'
            )
        # Circle
        step_svgs.append(f'  <circle cx="{cx}" cy="480" r="36" fill="{_esc(primary)}"/>')
        step_svgs.append(
            f'  <text x="{cx}" y="488" text-anchor="middle"'
            f' font-family="{_esc(font_title)}, sans-serif" font-size="28" font-weight="bold" fill="#FFFFFF">{i+1}</text>'
        )
        # Step text (wrap at ~20 chars)
        words = s[:50]
        step_svgs.append(
            f'  <text x="{cx}" y="560" text-anchor="middle"'
            f' font-family="{_esc(font_body)}, Microsoft YaHei, sans-serif"'
            f' font-size="24" fill="{_esc(text)}">{_esc(words)}</text>'
        )

    return f'''<svg viewBox="0 0 1280 960" xmlns="http://www.w3.org/2000/svg">
  <rect width="1280" height="960" fill="{_esc(bg)}"/>
  <rect x="0" y="0" width="1280" height="8" fill="{_esc(primary)}"/>
  <rect x="0" y="8" width="1280" height="140" fill="{_esc(primary)}"/>
  <text x="64" y="100" font-family="{_esc(font_title)}, Microsoft YaHei, sans-serif"
        font-size="44" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>
  <!-- Timeline track -->
  <rect x="{start_x + 36}" y="477" width="{1100 - 72}" height="6" rx="3" fill="{_esc(primary)}" opacity="0.2"/>
{chr(10).join(step_svgs)}
  <rect x="0" y="900" width="1280" height="4" fill="{_esc(accent)}" opacity="0.6"/>
</svg>'''


LAYOUT_BUILDERS = {
    "bullets": _build_bullets_svg,
    "two_col": _build_two_col_svg,
    "stats": _build_stats_svg,
    "timeline": _build_timeline_svg,
    "quote": _build_quote_svg,
    "section_break": _build_section_break_svg,
}


def _render_slides_fallback(outline: dict, design: dict) -> list[str]:
    """Generate SVG pages using the built-in template builders (no LLM)."""
    svgs: list[str] = []

    # Title slide
    svgs.append(_build_title_svg(
        outline.get("title", "Presentation"),
        outline.get("subtitle", ""),
        design,
    ))

    # Content slides
    for i, slide in enumerate(outline.get("slides", [])):
        layout = slide.get("layout", "bullets")
        builder = LAYOUT_BUILDERS.get(layout, _build_bullets_svg)
        svgs.append(builder(slide, design, i))

    return svgs


# ─────────────────────────────────────────────────────────────────────────
#  SVG → PPTX  (via ppt-master)
# ─────────────────────────────────────────────────────────────────────────

def _svgs_to_pptx(svg_pages: list[str], notes: dict[str, str] | None = None) -> str:
    """Write SVG files to a temp dir and call ppt-master to export PPTX."""

    if not _PPTMASTER_AVAILABLE:
        raise RuntimeError(
            "ppt-master skill not found. SVG import requires the optional ppt-master "
            "skill; the built-in python-pptx renderer (generate_pptx_from_outline) is "
            "available without it."
        )

    try:
        from svg_to_pptx import create_pptx_with_native_svg
    except ImportError as exc:
        raise RuntimeError(f"Failed to import svg_to_pptx from ppt-master: {exc}") from exc

    with tempfile.TemporaryDirectory(prefix="researchpal_ppt_") as tmp:
        tmp_path = Path(tmp)
        svg_dir = tmp_path / "svg_output"
        svg_dir.mkdir()

        svg_files: list[Path] = []
        for i, svg_code in enumerate(svg_pages):
            svg_path = svg_dir / f"slide_{i+1:02d}.svg"
            svg_path.write_text(svg_code, encoding="utf-8")
            svg_files.append(svg_path)

        os.makedirs(settings.OUTPUT_DIR, exist_ok=True)
        output_path = Path(settings.OUTPUT_DIR) / f"presentation_{uuid.uuid4().hex[:8]}.pptx"

        success = create_pptx_with_native_svg(
            svg_files=svg_files,
            output_path=output_path,
            canvas_format=CANVAS_FORMAT,
            verbose=False,
            transition="fade",
            notes=notes or {},
        )

        if not success:
            raise RuntimeError("ppt-master svg_to_pptx conversion returned failure")

        logger.info(f"ppt-master PPTX saved to {output_path}")
        return str(output_path)


# ─────────────────────────────────────────────────────────────────────────
#  Self-contained python-pptx renderer (no ppt-master dependency)
# ─────────────────────────────────────────────────────────────────────────

def _rgb(hex_color: str) -> RGBColor:
    h = (hex_color or "#1A56E8").lstrip("#")
    if len(h) != 6:
        h = "1A56E8"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _txt_on(bg_hex: str) -> RGBColor:
    return _rgb(_contrasting_text(bg_hex, "#FFFFFF", "#1F2937"))


def _add_bullets(slide, bullets, left, top, width, font_body, text, accent):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = "●  " + str(b)
        run.font.size = Pt(22)
        run.font.color.rgb = text
        run.font.name = font_body


def _outline_to_pptx(outline: dict, design: dict, notes: dict | None = None) -> str:
    """Self-contained PPTX renderer using python-pptx (no ppt-master dependency).

    Used as a graceful fallback when the optional ppt-master skill is not present,
    so PPT export always works out of the box (e.g. after a fresh git clone).
    """
    bg = _rgb(design.get("bg", "#FFFFFF"))
    primary = _rgb(design.get("primary", "#1A56E8"))
    accent = _rgb(design.get("accent", "#F59E0B"))
    text = _rgb(design.get("text", "#1F2937"))
    font_title = design.get("font_title", "Arial")
    font_body = design.get("font_body", "Arial")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width

    # ── Title slide ──
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    tb = s.shapes.add_textbox(Inches(1), Inches(2.9), Inches(11.33), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = outline.get("title", "Presentation")
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.name = font_title
    r.font.color.rgb = _txt_on(design.get("bg", "#1A56E8"))
    sub = s.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.33), Inches(1))
    tf2 = sub.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = outline.get("subtitle", "")
    r2.font.size = Pt(24)
    r2.font.name = font_body
    r2.font.color.rgb = _txt_on(design.get("bg", "#1A56E8"))

    slides = outline.get("slides", [])
    for i, slide in enumerate(slides):
        layout = slide.get("layout", "bullets")
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg

        # header bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()
        bar.shadow.inherit = False
        ht = s.shapes.add_textbox(Inches(0.5), Inches(0.18), W - Inches(1), Inches(0.85))
        hp = ht.text_frame.paragraphs[0]
        hr = hp.add_run()
        hr.text = slide.get("title", "")
        hr.font.size = Pt(30)
        hr.font.bold = True
        hr.font.name = font_title
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if notes and notes.get(f"slide_{i+2:02d}"):
            s.notes_slide.notes_text_frame.text = notes[f"slide_{i+2:02d}"]

        if layout == "two_col":
            bullets = slide.get("bullets", [])
            mid = len(bullets) // 2 or 1
            _add_bullets(s, bullets[:mid], 0.6, 1.6, 5.8, font_body, text, accent)
            _add_bullets(s, bullets[mid:], 6.9, 1.6, 5.8, font_body, text, accent)
        elif layout == "stats":
            cards = slide.get("bullets", [])[:4]
            n = len(cards)
            cw, gap = 2.7, 0.4
            total = n * cw + (n - 1) * gap
            start = (13.333 - total) / 2
            for j, c in enumerate(cards):
                cx = start + j * (cw + gap)
                card = s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.8), Inches(cw), Inches(2.6)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = primary
                card.line.fill.background()
                card.shadow.inherit = False
                ct = card.text_frame
                ct.word_wrap = True
                ct.vertical_anchor = MSO_ANCHOR.MIDDLE
                cp = ct.paragraphs[0]
                cp.alignment = PP_ALIGN.CENTER
                cr = cp.add_run()
                cr.text = str(c)
                cr.font.size = Pt(20)
                cr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cr.font.name = font_body
        elif layout == "timeline":
            steps = slide.get("bullets", [])[:5]
            n = len(steps)
            if n == 0:
                _add_bullets(s, slide.get("bullets", []), 0.6, 1.6, 12.1, font_body, text, accent)
            else:
                step_w = 11.0 / n
                for j, st in enumerate(steps):
                    cx = 0.6 + j * step_w + step_w / 2
                    circ = s.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(cx - 0.4), Inches(2.6), Inches(0.8), Inches(0.8)
                    )
                    circ.fill.solid()
                    circ.fill.fore_color.rgb = primary
                    circ.line.fill.background()
                    circ.shadow.inherit = False
                    ctf = circ.text_frame
                    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cp = ctf.paragraphs[0]
                    cp.alignment = PP_ALIGN.CENTER
                    cr = cp.add_run()
                    cr.text = str(j + 1)
                    cr.font.size = Pt(24)
                    cr.font.bold = True
                    cr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    tt = s.shapes.add_textbox(Inches(cx - 1.3), Inches(3.6), Inches(2.6), Inches(1.6))
                    tp = tt.text_frame
                    tp.word_wrap = True
                    tr = tp.paragraphs[0]
                    tr.alignment = PP_ALIGN.CENTER
                    tr2 = tr.add_run()
                    tr2.text = str(st)[:50]
                    tr2.font.size = Pt(16)
                    tr2.font.color.rgb = text
                    tr2.font.name = font_body
        elif layout in ("quote", "section_break"):
            s.background.fill.fore_color.rgb = primary
            big = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.33), Inches(2.4))
            bf = big.text_frame
            bf.word_wrap = True
            bf.vertical_anchor = MSO_ANCHOR.MIDDLE
            bp = bf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = (slide.get("bullets", [""])[0] if slide.get("bullets") else "")
            br.font.size = Pt(32)
            br.font.italic = True
            br.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            br.font.name = font_title
            if len(slide.get("bullets", [])) > 1:
                src = s.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.33), Inches(1))
                sp = src.text_frame.paragraphs[0]
                sp.alignment = PP_ALIGN.CENTER
                sr = sp.add_run()
                sr.text = slide["bullets"][1]
                sr.font.size = Pt(22)
                sr.font.color.rgb = accent
                sr.font.name = font_body
        else:  # bullets
            _add_bullets(s, slide.get("bullets", []), 0.8, 1.6, 11.7, font_body, text, accent)

    os.makedirs(settings.OUTPUT_DIR, exist_ok=True)
    output_path = Path(settings.OUTPUT_DIR) / f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(str(output_path))
    logger.info(f"python-pptx PPTX saved to {output_path}")
    return str(output_path)


# ─────────────────────────────────────────────────────────────────────────
#  PUBLIC API
# ─────────────────────────────────────────────────────────────────────────

def generate_pptx_from_outline(outline: dict, design: dict) -> str:
    """
    Generate a PPTX from a structured outline + design spec.

    Priority: built-in SVG builders → ppt-master svg_to_pptx (if the optional
    skill is present); otherwise falls back to the self-contained python-pptx
    renderer so PPT export always works out of the box.
    """
    # Build speaker notes dict
    notes: dict[str, str] = {}
    slides = outline.get("slides", [])
    for i, slide in enumerate(slides):
        n = slide.get("notes", "")
        if n:
            # slide index 0 = title slide, so content starts at 1
            notes[f"slide_{i+2:02d}"] = n

    if _PPTMASTER_AVAILABLE:
        try:
            svg_pages = _render_slides_fallback(outline, design)
            return _svgs_to_pptx(svg_pages, notes)
        except Exception as e:  # pragma: no cover - defensive fallback
            logger.warning(f"ppt-master export failed, falling back to python-pptx: {e}")

    return _outline_to_pptx(outline, design, notes)


def generate_pptx_from_svgs(svg_pages: list[str], notes: dict[str, str] | None = None) -> str:
    """
    Generate a PPTX from pre-rendered SVG strings (e.g., LLM-generated).
    Requires the optional ppt-master skill.
    """
    return _svgs_to_pptx(svg_pages, notes)
