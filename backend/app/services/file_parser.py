import os
import json
import logging
import chardet
from typing import Optional

logger = logging.getLogger(__name__)


def parse_pdf(file_path: str) -> dict:
    import pdfplumber

    text_parts = []
    pages_count = 0
    metadata = {}
    try:
        with pdfplumber.open(file_path) as pdf:
            pages_count = len(pdf.pages)
            metadata = pdf.metadata or {}
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        raise ValueError(f"Failed to parse PDF: {e}")

    return {
        "text": "\n\n".join(text_parts),
        "pages_count": pages_count,
        "metadata": {k: str(v) for k, v in metadata.items() if v},
    }


def parse_docx(file_path: str) -> dict:
    from docx import Document

    doc = Document(file_path)
    parts = []

    # Extract paragraphs
    para_count = 0
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
            para_count += 1

    # Extract tables
    table_count = 0
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append(" | ".join(cells))
        if table_rows:
            parts.append("[表格]\n" + "\n".join(table_rows))
            table_count += 1

    return {
        "text": "\n\n".join(parts),
        "paragraphs_count": para_count,
        "tables_count": table_count,
    }


def parse_pptx(file_path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        slide_parts = [f"--- 第 {slide_count} 页 ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                table = shape.table
                table_rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    slide_parts.append("[表格]\n" + "\n".join(table_rows))
        if len(slide_parts) > 1:
            slides.append("\n".join(slide_parts))
        else:
            slides.append(f"--- 第 {slide_count} 页 --- [空白页或仅图片]")

    return {
        "text": "\n\n".join(slides),
        "slides_count": slide_count,
    }


def parse_csv(file_path: str) -> dict:
    import pandas as pd

    df = pd.read_csv(file_path)
    columns = list(df.columns)
    row_count = len(df)
    preview = df.head(10).to_dict(orient="records")
    return {
        "columns": columns,
        "row_count": row_count,
        "preview": preview,
        "dataframe_json": df.to_json(orient="records", force_ascii=False),
    }


def parse_excel(file_path: str) -> dict:
    import pandas as pd

    xls = pd.ExcelFile(file_path)
    sheets = {}
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        sheets[sheet_name] = {
            "columns": list(df.columns),
            "row_count": len(df),
            "preview": df.head(10).to_dict(orient="records"),
        }
    return {
        "sheets": list(xls.sheet_names),
        "dataframes": sheets,
    }


def parse_text(file_path: str) -> dict:
    with open(file_path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8")
    try:
        text = raw.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = raw.decode("utf-8", errors="replace")
        encoding = "utf-8 (fallback)"
    return {
        "text": text,
        "encoding": encoding,
    }


def extract_text(file_id: str, file_path: str, file_type: str) -> dict:
    file_type = file_type.lower()
    if file_type == "pdf":
        result = parse_pdf(file_path)
    elif file_type == "docx":
        result = parse_docx(file_path)
    elif file_type == "doc":
        # Old .doc format: python-docx can't read it; try as fallback
        raise ValueError("旧版 .doc 格式不支持，请另存为 .docx 格式后重新上传")
    elif file_type == "pptx":
        result = parse_pptx(file_path)
    elif file_type == "csv":
        result = parse_csv(file_path)
    elif file_type in ("xlsx", "xls"):
        result = parse_excel(file_path)
    elif file_type in ("txt", "md", "py", "json", "log"):
        result = parse_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    result["file_id"] = file_id
    result["file_type"] = file_type
    return result
